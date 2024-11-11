import os
from pptx import Presentation

# Define the directory containing the PowerPoint files
# directory = "c:\checkouts\river_course\lectures\day4\"
directory = os.path.dirname(os.path.abspath(__file__))

# Text to add or update on the first slide
new_title_text = "An awesome presentation"
new_subtitle_text = "Really"
old_author = "Victor"
new_author = "Mr. 42\nSobek"
new_image = "Sobek.png" # Path to the new image file

# Loop through all files in the directory   
for filename in os.listdir(directory):
    if filename.endswith(".pptx"):
        ppt_path = os.path.join(directory, filename)
        prs = Presentation(ppt_path)
        
        # Access the first slide
        first_slide = prs.slides[0]
        
        # Modify the title
        if first_slide.shapes.title:
            first_slide.shapes.title.text = new_title_text
        
        for shape in first_slide.shapes:
            # Modify subtitle 
            if shape.has_text_frame and "subtitle" in shape.name.lower():
                shape.text = new_subtitle_text
            # Modify author    
            if shape.has_text_frame and old_author in shape.text:
                nlines=len(shape.text_frame.paragraphs)
                alignment_old = shape.text_frame.paragraphs[0].alignment
                shape.text = new_author
                for k in range(nlines):
                    shape.text_frame.paragraphs[k].alignment = alignment_old
            # Modify image
            if not shape.has_text_frame: #then it must be the picture
            #if shape.shape_type == 13:  # 13 is the shape type for a picture, but does not seem to be valid in this presentation?
                #shape.image.from_file = new_image  
                # get part and rId from shape we need to change
                slide_part, rId = shape.part, shape._element.blip_rId
                image_part = slide_part.related_part(rId)

                # overwrite old blob info with new blob info
                new_pptx_img = pptx.parts.image.Image.from_file(new_image)
                image_part.blob = new_pptx_img._blob

        # Save the presentation
        prs.save(os.path.join(directory, f"updated_{filename}"))

print("First slide updated for all PowerPoint files.")